import ffmpeg
import tempfile
import os
import aiosqlite
import asyncio
import logging
import io
import base64
import random
import httpx
from datetime import datetime, timedelta
from aiogram import Bot, Dispatcher, F
from aiogram.types import Message, CallbackQuery, BufferedInputFile
from aiogram.types import Message, CallbackQuery, InlineKeyboardMarkup, InlineKeyboardButton, BufferedInputFile
from aiogram.filters import CommandStart, Command
from aiogram.fsm.storage.memory import MemoryStorage
from openai import AsyncOpenAI
from pptx import Presentation
from docx import Document as DocxDocument
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import matplotlib
from config import BOT_TOKEN, CLAUDE_API_KEY, CLAUDE_API_KEY_2, ADMIN_ID, PLANS, UNSPLASH_KEY
from logger_files import log_message, log_photo, log_user, log_voice, log_circle, init_logs
matplotlib.use('Agg')

from database import (
    init_db, get_user, create_user, increment_messages,
    can_send_message, save_message, get_history,
    clear_history, get_stats, upgrade_user,
    ACHIEVEMENTS, give_achievement, get_user_achievements,
    get_user_stats, update_user_stats, get_top_users,
    get_weekly_top_user, add_bonus_messages, update_streak,
    add_referral, get_referral_stats
)
from keyboards import main_menu, plans_keyboard

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
matplotlib.use('Agg')

UPCOMING_UPDATE = """
🚀 <b>Будущее обновление</b>

✨ Что планируется:
• Telegram подарки за достижения
• Улучшение дизана презентаций

📅 Ожидаемая дата: скоро
"""

WEEKLY_BONUS = 50
CHANNEL_ID = "@Neirosupport1"
CHANNEL_BONUS = 10

bot = Bot(token=BOT_TOKEN)
dp = Dispatcher(storage=MemoryStorage())
@dp.message(F.chat.type.in_({"channel", "supergroup", "group"}))
async def ignore_groups(message: Message):
    return
claude = AsyncOpenAI(
    api_key=CLAUDE_API_KEY,
    base_url="https://api.claudehub.fun/v1"
)
CLAUDE_KEYS = [CLAUDE_API_KEY, CLAUDE_API_KEY_2]
current_key_index = 0

def get_claude_client():
    global current_key_index
    key = CLAUDE_KEYS[current_key_index % len(CLAUDE_KEYS)]
    current_key_index += 1
    return AsyncOpenAI(api_key=key, base_url="https://api.claudehub.fun/v1")

# ===================== ДОСТИЖЕНИЯ =====================
async def check_and_give(user_id: int, achievement_id: str):
    given = await give_achievement(user_id, achievement_id)
    if given:
        a = ACHIEVEMENTS[achievement_id]
        
        if a.get("reward_type") == "gift":
            reward_text = f"\n🎁 Награда: подарок <b>{a['reward']} ⭐</b>"
            # Уведомление админу
            user = await get_user(user_id)
            username = f"@{user['username']}" if user and user['username'] else f"ID: {user_id}"
            try:
                await bot.send_message(
                    ADMIN_ID,
                    f"🎁 <b>Нужно выдать подарок!</b>\n\n"
                    f"👤 Пользователь: {username}\n"
                    f"🏅 Достижение: {a['emoji']} {a['name']}\n"
                    f"⭐ Стоимость: <b>{a['reward']} звёзд</b>",
                    parse_mode="HTML"
                )
            except Exception:
                pass
        else:
            reward_text = f"\n🎁 Награда: <b>+{a['reward']} запросов</b>" if a.get("reward", 0) > 0 else ""

        try:
            await bot.send_message(
                user_id,
                f"🎉 <b>Новое достижение!</b>\n\n"
                f"{a['emoji']} <b>{a['name']}</b>\n"
                f"<i>{a['desc']}</i>"
                f"{reward_text}",
                parse_mode="HTML"
            )
        except Exception:
            pass

async def weekly_top_scheduler():
    while True:
        now = datetime.now()
        days_until_monday = (7 - now.weekday()) % 7 or 7
        next_monday = (now + timedelta(days=days_until_monday)).replace(
            hour=9, minute=0, second=0, microsecond=0
        )
        wait_seconds = (next_monday - now).total_seconds()
        await asyncio.sleep(wait_seconds)

        try:
            top = await get_weekly_top_user()
            if not top:
                continue

            name = top["username"] or f"user{top['user_id']}"
            await check_and_give(top["user_id"], "week_champion")

            try:
                await bot.send_message(
                    top["user_id"],
                    f"🏆 <b>Поздравляем!</b>\n\n"
                    f"Ты лидер этой недели!\n"
                    f"💬 Сообщений: {top['messages_used']}\n"
                    f"📊 Презентаций: {top['pres']}\n\n"
                    f"⏳ Скоро получишь бонусные запросы от администратора!",
                    parse_mode="HTML"
                )
            except Exception:
                pass

            await bot.send_message(
                ADMIN_ID,
                f"📊 <b>Итоги недели!</b>\n\n"
                f"🏆 Победитель: <b>{name}</b>\n"
                f"🆔 ID: <code>{top['user_id']}</code>\n"
                f"💬 Сообщений: {top['messages_used']}\n"
                f"📊 Презентаций: {top['pres']}\n"
                f"📸 Фото: {top['photos']}\n\n"
                f"Выдать бонус {WEEKLY_BONUS} запросов:\n"
                f"<code>/bonus {top['user_id']} {WEEKLY_BONUS}</code>",
                parse_mode="HTML"
            )
        except Exception as e:
            logger.error(f"Weekly scheduler error: {e}")

# ===================== ВСПОМОГАТЕЛЬНЫЕ =====================
async def claude_request(**kwargs):
    try:
        client = get_claude_client()
        response = await client.chat.completions.create(**kwargs)
        return response
    except Exception as e:
        logger.error(f"Claude request error: {e}")
        raise

THEMES = [
    {
        "name": "violet",
        "bg_dark": RGBColor(15, 12, 41),
        "bg_light": RGBColor(240, 238, 255),
        "accent": RGBColor(99, 102, 241),
        "accent2": RGBColor(16, 185, 129),
        "accent3": RGBColor(245, 158, 11),
        "white": RGBColor(255, 255, 255),
        "text_dark": RGBColor(15, 12, 41),
        "text_light": RGBColor(200, 210, 255),
    },
    {
        "name": "ocean",
        "bg_dark": RGBColor(10, 25, 47),
        "bg_light": RGBColor(230, 245, 255),
        "accent": RGBColor(0, 180, 216),
        "accent2": RGBColor(0, 245, 212),
        "accent3": RGBColor(255, 107, 107),
        "white": RGBColor(255, 255, 255),
        "text_dark": RGBColor(10, 25, 47),
        "text_light": RGBColor(180, 230, 255),
    },
    {
        "name": "fire",
        "bg_dark": RGBColor(26, 10, 10),
        "bg_light": RGBColor(255, 245, 240),
        "accent": RGBColor(239, 68, 68),
        "accent2": RGBColor(251, 146, 60),
        "accent3": RGBColor(250, 204, 21),
        "white": RGBColor(255, 255, 255),
        "text_dark": RGBColor(26, 10, 10),
        "text_light": RGBColor(255, 200, 180),
    },
    {
        "name": "forest",
        "bg_dark": RGBColor(10, 30, 20),
        "bg_light": RGBColor(235, 255, 245),
        "accent": RGBColor(16, 185, 129),
        "accent2": RGBColor(52, 211, 153),
        "accent3": RGBColor(99, 102, 241),
        "white": RGBColor(255, 255, 255),
        "text_dark": RGBColor(10, 30, 20),
        "text_light": RGBColor(180, 255, 220),
    },
    {
        "name": "luxury",
        "bg_dark": RGBColor(20, 15, 5),
        "bg_light": RGBColor(255, 252, 240),
        "accent": RGBColor(212, 175, 55),
        "accent2": RGBColor(255, 215, 0),
        "accent3": RGBColor(192, 192, 192),
        "white": RGBColor(255, 255, 255),
        "text_dark": RGBColor(20, 15, 5),
        "text_light": RGBColor(255, 240, 180),
    },
]

def get_random_theme():
    return random.choice(THEMES)

def rgb_to_hex(color):
    if isinstance(color, RGBColor):
        return '#{:02x}{:02x}{:02x}'.format(color[0], color[1], color[2])
    return '#{:02x}{:02x}{:02x}'.format(color.red, color.green, color.blue)

def get_chart_colors(theme):
    return [
        rgb_to_hex(theme["accent"]),
        rgb_to_hex(theme["accent2"]),
        rgb_to_hex(theme["accent3"]),
        '#8b5cf6', '#ec4899', '#06b6d4', '#84cc16'
    ]

async def get_wiki_image(query):
    headers = {"User-Agent": "Mozilla/5.0 (compatible; TelegramBot/1.0; +https://t.me/bot)"}
    for lang in ["ru", "en"]:
        try:
            async with httpx.AsyncClient(timeout=10, headers=headers) as client:
                search_params = {"action": "query", "format": "json", "list": "search", "srsearch": query, "srlimit": 1}
                r = await client.get(f"https://{lang}.wikipedia.org/w/api.php", params=search_params)
                data = r.json()
                results = data.get("query", {}).get("search", [])
                if not results:
                    continue
                title = results[0]["title"]
                img_params = {"action": "query", "format": "json", "prop": "pageimages", "titles": title, "pithumbsize": 500, "pilimit": 1}
                r2 = await client.get(f"https://{lang}.wikipedia.org/w/api.php", params=img_params)
                data2 = r2.json()
                pages = data2.get("query", {}).get("pages", {})
                for page in pages.values():
                    img_url = page.get("thumbnail", {}).get("source")
                    if img_url:
                        img_r = await client.get(img_url)
                        return img_r.content
        except Exception as e:
            logger.error(f"Wiki image error ({lang}): {e}")
    return None

async def get_unsplash_image(query):
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            r = await client.get(
                "https://api.unsplash.com/search/photos",
                params={"query": query, "per_page": 1, "orientation": "landscape"},
                headers={"Authorization": f"Client-ID {UNSPLASH_KEY}"}
            )
            data = r.json()
            results = data.get("results", [])
            if results:
                img_url = results[0]["urls"]["small"]
                img_r = await client.get(img_url)
                return img_r.content
    except Exception as e:
        logger.error(f"Unsplash error: {e}")
    return None

# ===================== СЛАЙДЫ =====================
def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def add_colored_rectangle(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rectangle(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def create_chart_image(chart_type, labels, values, title, theme):
    colors_chart = get_chart_colors(theme)
    bg_hex = rgb_to_hex(theme["bg_dark"])
    accent_hex = rgb_to_hex(theme["accent"])
    fig, ax = plt.subplots(figsize=(6, 4), facecolor=bg_hex)
    ax.set_facecolor(bg_hex)
    if chart_type == "bar":
        bars = ax.bar(labels, values, color=colors_chart[:len(labels)], edgecolor='white', linewidth=0.5, width=0.6)
        for bar, val in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(values)*0.02,
                   str(int(val)), ha='center', color='white', fontsize=10, fontweight='bold')
        ax.set_ylim(0, max(values) * 1.2)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        for spine in ['bottom', 'left']:
            ax.spines[spine].set_edgecolor(accent_hex)
    elif chart_type == "pie":
        wedges, texts, autotexts = ax.pie(
            values, labels=labels, colors=colors_chart[:len(labels)],
            autopct='%1.1f%%', textprops={'color': 'white', 'fontsize': 9},
            wedgeprops={'edgecolor': bg_hex, 'linewidth': 2}, startangle=90
        )
        for autotext in autotexts:
            autotext.set_fontweight('bold')
    elif chart_type == "line":
        ax.plot(labels, values, color=accent_hex, linewidth=3, marker='o', markersize=10,
               markerfacecolor=rgb_to_hex(theme["accent2"]), markeredgecolor='white', markeredgewidth=2)
        ax.fill_between(range(len(labels)), values, alpha=0.2, color=accent_hex)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        for spine in ['bottom', 'left']:
            ax.spines[spine].set_edgecolor(accent_hex)
    elif chart_type == "horizontal_bar":
        bars = ax.barh(labels, values, color=colors_chart[:len(labels)], edgecolor='white', linewidth=0.5)
        for bar, val in zip(bars, values):
            ax.text(val + max(values)*0.02, bar.get_y() + bar.get_height()/2,
                   str(int(val)), va='center', color='white', fontsize=10, fontweight='bold')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        for spine in ['bottom', 'left']:
            ax.spines[spine].set_edgecolor(accent_hex)
    ax.set_title(title, color='white', fontsize=12, pad=15, fontweight='bold')
    ax.tick_params(colors='white', labelsize=9)
    ax.grid(axis='y' if chart_type in ['bar', 'line'] else 'x', alpha=0.2, color='white', linestyle='--')
    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', facecolor=bg_hex, dpi=150)
    plt.close()
    buf.seek(0)
    return buf

def create_title_slide(prs, title, subtitle, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, theme["bg_dark"])
    add_colored_rectangle(slide, 0, 0, 10, 1.0, theme["accent"])
    add_colored_rectangle(slide, 0, 6.5, 10, 1.0, theme["accent2"])
    add_colored_rectangle(slide, 0, 0, 0.3, 7.5, theme["accent3"])
    add_rounded_rectangle(slide, 7.5, 1.5, 2.0, 2.0, theme["accent"])
    add_rounded_rectangle(slide, 8.0, 2.0, 1.0, 1.0, theme["accent2"])
    add_text_box(slide, title, 0.6, 1.5, 8.5, 1.8, font_size=38, bold=True, color=theme["white"], align=PP_ALIGN.LEFT)
    add_text_box(slide, subtitle, 0.6, 4.0, 8.5, 0.8, font_size=20, color=theme["text_light"], italic=True)
    add_text_box(slide, "✦ Powered by AI", 0.6, 5.2, 4, 0.5, font_size=13, color=theme["accent2"])
    return slide

def create_content_slide(prs, title, content, slide_num, theme, has_chart=False, chart_data=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if slide_num % 2 == 0:
        set_slide_background(slide, theme["bg_dark"])
        text_color = theme["white"]
    else:
        set_slide_background(slide, theme["bg_light"])
        text_color = theme["text_dark"]
    add_colored_rectangle(slide, 0, 0, 10, 1.3, theme["accent"])
    add_colored_rectangle(slide, 0, 1.3, 10, 0.07, theme["accent2"])
    add_rounded_rectangle(slide, 8.5, 0.15, 0.9, 0.9, theme["accent2"])
    add_text_box(slide, f"{slide_num:02d}", 8.52, 0.18, 0.86, 0.8, font_size=22, bold=True, color=theme["white"], align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 0.3, 0.15, 8.0, 1.0, font_size=24, bold=True, color=theme["white"])
    add_colored_rectangle(slide, 0, 1.4, 0.12, 6.1, theme["accent3"])
    if has_chart and chart_data:
        add_text_box(slide, content, 0.3, 1.5, 4.6, 4.8, font_size=14, color=text_color)
        try:
            chart_img = create_chart_image(chart_data["type"], chart_data["labels"], chart_data["values"], chart_data["title"], theme)
            slide.shapes.add_picture(chart_img, Inches(5.0), Inches(1.5), Inches(4.7), Inches(3.8))
        except Exception as e:
            logger.error(f"Chart error: {e}")
    else:
        lines = [l for l in content.split("\n") if l.strip()]
        y_pos = 1.55
        for i, line in enumerate(lines[:7]):
            marker_color = theme["accent"] if i % 2 == 0 else theme["accent2"]
            add_rounded_rectangle(slide, 0.25, y_pos + 0.08, 0.25, 0.28, marker_color)
            add_text_box(slide, line.strip().lstrip("-▸•").strip(), 0.65, y_pos, 6.3, 0.55, font_size=15, color=text_color)
            y_pos += 0.65
    add_colored_rectangle(slide, 0, 7.1, 10, 0.4, theme["accent"])
    add_text_box(slide, title, 0.3, 7.12, 8, 0.3, font_size=10, color=theme["text_light"])
    return slide

def create_divider_slide(prs, text, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, theme["accent"])
    add_colored_rectangle(slide, 0, 2.8, 10, 0.1, theme["white"])
    add_colored_rectangle(slide, 0, 4.5, 10, 0.1, theme["accent2"])
    add_text_box(slide, "◆◆◆", 0, 2.0, 10, 0.6, font_size=20, color=theme["white"], align=PP_ALIGN.CENTER)
    add_text_box(slide, text, 0, 3.0, 10, 1.2, font_size=34, bold=True, color=theme["white"], align=PP_ALIGN.CENTER)
    return slide

def create_summary_slide(prs, title, points, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, theme["bg_dark"])
    add_colored_rectangle(slide, 0, 0, 10, 1.5, theme["accent2"])
    add_colored_rectangle(slide, 0, 6.8, 10, 0.7, theme["accent"])
    add_colored_rectangle(slide, 4.5, 1.6, 0.1, 5.5, theme["accent3"])
    add_text_box(slide, "ИТОГИ И ВЫВОДЫ", 0, 0.1, 10, 0.7, font_size=14, bold=True, color=theme["white"], align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 0, 0.7, 10, 0.8, font_size=28, bold=True, color=theme["white"], align=PP_ALIGN.CENTER)
    y = 1.8
    colors_cycle = [theme["accent"], theme["accent2"], theme["accent3"]]
    for i, point in enumerate(points[:5], 1):
        c = colors_cycle[i % len(colors_cycle)]
        add_rounded_rectangle(slide, 0.4, y, 0.6, 0.5, c)
        add_text_box(slide, str(i), 0.42, y + 0.02, 0.56, 0.45, font_size=18, bold=True, color=theme["white"], align=PP_ALIGN.CENTER)
        add_text_box(slide, point, 1.2, y, 8.4, 0.55, font_size=14, color=theme["white"])
        y += 0.78
    add_text_box(slide, "🚀 Спасибо за внимание!", 0, 6.85, 10, 0.4, font_size=16, bold=True, color=theme["white"], align=PP_ALIGN.CENTER)
    return slide

# ===================== КОМАНДЫ =====================
@dp.message(CommandStart())
async def start(message: Message):
    user = message.from_user
    is_new = await get_user(user.id) is None
    await create_user(user.id, user.username or user.first_name)

    args = message.text.split()
    if is_new and len(args) > 1 and args[1].startswith("ref_"):
        try:
            referrer_id = int(args[1].replace("ref_", ""))
            if referrer_id != user.id:
                success = await add_referral(referrer_id, user.id)
                if success:
                    await bot.send_message(
                        referrer_id,
                        f"🎉 <b>Новый реферал!</b>\n\n"
                        f"По твоей ссылке зарегистрировался новый пользователь.\n"
                        f"Тебе начислено <b>+15 запросов!</b> 🎁",
                        parse_mode="HTML"
                    )
                    await message.answer(
                        f"🎁 <b>Бонус за приглашение!</b>\n\n"
                        f"Тебе начислено <b>+15 запросов</b> по реферальной ссылке!",
                        parse_mode="HTML"
                    )
        except:
            pass

    await log_user(user.id, user.username or "", user.first_name, user.last_name or "")
    await check_and_give(user.id, "first_step")
    stats = await get_user_stats(user.id)
    used_cmds = set(stats["commands_used"].split(",")) if stats["commands_used"] else set()
    used_cmds.add("start")
    await update_user_stats(user.id, commands_used=",".join(used_cmds))

    await message.answer(
        f"👋 Привет, {user.first_name}!\n\n"
        f"Я AI-ассистент на базе Claude.\n\n"
        f"🆓 У тебя 10 бесплатных сообщений в месяц.\n"
        f"Чтобы получить больше — смотри /plans\n\n"
        f"📊 Могу создать презентацию:\n"
        f"<code>создай презентацию [тема]</code>",
        reply_markup=main_menu(),
        parse_mode="HTML"
    )

@dp.message(F.text == "💬 Новый чат")
async def new_chat(message: Message):
    await clear_history(message.from_user.id)
    await message.answer("💬 Новый чат начат! Пиши что хочешь.")

@dp.message(F.text == "👤 Профиль")
@dp.message(Command("profile"))
async def profile(message: Message):
    user_id = message.from_user.id
    user = await get_user(user_id)
    if not user:
        await start(message)
        return

    user_achievements = await get_user_achievements(user_id)
    earned_ids = {r["achievement_id"] for r in user_achievements}
    earned_count = len(earned_ids)
    total_count = len(ACHIEVEMENTS)

    plan_name = PLANS.get(user["plan"], {}).get("name", user["plan"])
    remaining = user["messages_limit"] - user["messages_used"]

    plan_features = {
        "free":  "• 10 сообщений в месяц\n• 3 презентации\n• Базовый анализ фото",
        "basic": "• 100 сообщений в месяц\n• 20 презентаций\n• Анализ фото\n• Приоритетный ответ",
        "pro":   "• 500 сообщений в месяц\n• Безлимит презентаций\n• Анализ фото\n• Быстрый ответ",
        "ultra": "• Безлимит сообщений\n• Безлимит презентаций\n• Все функции\n• Максимальный приоритет",
    }
    features = plan_features.get(user["plan"], "")

    earned_list = ""
    for aid in earned_ids:
        if aid in ACHIEVEMENTS:
            a = ACHIEVEMENTS[aid]
            earned_list += f"{a['emoji']} {a['name']}\n"

    await message.answer(
        f"👤 <b>Ваш профиль</b>\n\n"
        f"🗂 Тариф: <b>{plan_name}</b>\n"
        f"💬 Использовано: {user['messages_used']}/{user['messages_limit']}\n"
        f"✉️ Осталось: {remaining} сообщений\n"
        f"🔄 Сброс: {user['reset_date'][:10]}\n\n"
        f"📦 <b>Что включено:</b>\n{features}\n\n"
        f"🏆 <b>Достижения [{earned_count}/{total_count}]:</b>\n"
        f"{earned_list if earned_list else 'Пока нет достижений'}",
        parse_mode="HTML"
    )
    user_id = message.from_user.id
    user = await get_user(user_id)
    if not user:
        await start(message)
        return

    user_achievements = await get_user_achievements(user_id)
    earned = [ACHIEVEMENTS[r["achievement_id"]]["emoji"] + " " + ACHIEVEMENTS[r["achievement_id"]]["name"]
              for r in user_achievements if r["achievement_id"] in ACHIEVEMENTS]

    plan_name = PLANS.get(user["plan"], {}).get("name", user["plan"])
    remaining = user["messages_limit"] - user["messages_used"]
    ach_text = ", ".join(earned) if earned else "Пока нет"

    plan_features = {
        "free":  "• 10 сообщений в месяц на все(презентации включены в это колличество)\n• 3 презентации (\n• Базовый анализ фото",
        "basic": "• 100 сообщений в месяц на все(презентации включены в это колличество)\n• 20 презентаций(\n• Анализ фото\n• Приоритетный ответ",
        "pro":   "• 500 сообщений в месяц на все(презентации включены в это колличество)\n• Безлимит презентаций\n• Анализ фото\n• Быстрый ответ\n• Доступ к новым функциям",
        "ultra": "• Безлимит сообщений\n• Безлимит презентаций\n• Все функции\n• Максимальный приоритет",
    }
    features = plan_features.get(user["plan"], "")

    user_id = message.from_user.id
    user_achievements = await get_user_achievements(user_id)
    earned_ids = {row["achievement_id"] for row in user_achievements}

    categories = {}
    for aid, a in ACHIEVEMENTS.items():
        cat = a["cat"]
        if cat not in categories:
            categories[cat] = []
        if aid in earned_ids:
            categories[cat].append(f"{a['emoji']} <b>{a['name']}</b> — {a['desc']}")
        else:
            categories[cat].append(f"🔒 <i>{a['name']}</i>")

    text = f"🏆 <b>Достижения</b> [{len(earned_ids)}/{len(ACHIEVEMENTS)}]\n\n"
    for cat, items_aids in categories.items():
        text += f"<b>{cat}</b>\n"
        for aid, a in [(aid, ACHIEVEMENTS[aid]) for aid in ACHIEVEMENTS if ACHIEVEMENTS[aid]["cat"] == cat]:
            if aid in earned_ids:
                if a.get("reward_type") == "gift":
                    reward = f"⭐ {a['reward']} звёзд"
                elif a.get("reward", 0) > 0:
                    reward = f"+{a['reward']} запросов"
                else:
                    reward = ""
                reward_text = f" — <i>{reward}</i>" if reward else ""
                text += f"  ✅ {a['emoji']} <b>{a['name']}</b>{reward_text}\n"
            else:
                text += f"  🔒 {a['emoji']} {a['name']} — <i>{a['desc']}</i>\n"
        text += "\n"
    await message.answer(text, parse_mode="HTML")
@dp.message(Command("gifts"))
async def get_gifts(message: Message):
    if message.from_user.id != ADMIN_ID:
        return
    result = await bot.get_available_gifts()
    text = "🎁 <b>Доступные подарки:</b>\n\n"
    for gift in result.gifts:
        text += f"ID: <code>{gift.id}</code>\n"
        text += f"Цена: {gift.star_count} ⭐\n\n"
    await message.answer(text, parse_mode="HTML")
@dp.message(Command("testgift"))
async def test_gift(message: Message):
    if message.from_user.id != ADMIN_ID:
        return
    try:
        await bot.send_gift(
            user_id=message.from_user.id,
            gift_id="5170145012310081615"
        )
        await message.answer("✅ Подарок отправлен!")
    except Exception as e:
        await message.answer(f"❌ Ошибка: {e}")


@dp.message(Command("top"))
@dp.message(F.text == "📈 Топ")
async def show_top(message: Message):
    rows = await get_top_users(10)
    text = "📈 <b>Топ пользователей</b>\n\n"
    text += "👨‍💻 <b>#0 ZigelDev</b> — <i>Создатель этого безумия</i>\n   ∞ сообщений  ∞ презентаций  ∞ фото  🏆∞\n\n"
    text += "─────────────────\n\n"
    
    medals = ["🥇", "🥈", "🥉"]
    plan_labels = {
        "free":  "🆓",
        "basic": "⭐",
        "pro":   "🚀",
        "mega":  "👾"
    }

    rank = 0
    for row in rows:
        if row["user_id"] == ADMIN_ID:
            continue
        rank += 1
        medal = medals[rank-1] if rank <= 3 else f"{rank}."
        name = row["username"] or f"user{row['user_id']}"
        ach = f"🏆{row['achievements']}" if row["achievements"] else ""
        plan = plan_labels.get(dict(row).get("plan", "free"), "🆓")
        text += (
            f"{medal} <b>{name}</b> {plan}\n"
            f"   💬 {row['messages_used']}  📊 {row['pres']}  📸 {row['photos']}  {ach}\n\n"
        )

    await message.answer(text, parse_mode="HTML")

@dp.message(F.text == "💳 Тарифы")
@dp.message(Command("plans"))
async def show_plans(message: Message):
    text = "💳 <b>Тарифные планы</b>\n\n"
    for plan_id, plan in PLANS.items():
        price = f"{plan['price'] // 100}₽/мес" if plan["price"] > 0 else "Бесплатно"
        text += f"{plan['name']}\n"
        text += f"   📨 {plan['messages']} сообщений в месяц\n"
        text += f"   💰 {price}\n\n"
    await message.answer(text, parse_mode="HTML", reply_markup=plans_keyboard())

@dp.callback_query(F.data.startswith("buy_"))
async def buy_plan(callback: CallbackQuery):
    plan_id = callback.data.split("_")[1]
    plan = PLANS.get(plan_id)
    if not plan:
        await callback.answer("Тариф не найден")
        return
    user = callback.from_user
    price = plan['price'] // 100
    await callback.message.answer(
        f"💳 <b>Заявка на тариф {plan['name']}</b>\n\n"
        f"Сумма: <b>{price}₽</b>\n\n"
        f"📩 Для оплаты напиши администратору: @zigeli\n"
        f"После оплаты администратор активирует тариф вручную.",
        parse_mode="HTML"
    )
    await bot.send_message(
        ADMIN_ID,
        f"🔔 <b>Новая заявка на оплату!</b>\n\n"
        f"👤 Пользователь: {user.first_name} (@{user.username or 'нет'})\n"
        f"🆔 ID: <code>{user.id}</code>\n"
        f"📋 Тариф: {plan['name']}\n"
        f"💰 Сумма: {price}₽\n\n"
                f"Команда для активации:\n"
        f"<code>/give {user.id} {plan_id}</code>",
        parse_mode="HTML"
    )
    await callback.answer()

@dp.message(Command("give"))
async def give_plan(message: Message):
    if message.from_user.id != ADMIN_ID:
        return
    args = message.text.split()
    if len(args) != 3:
        await message.answer("❌ Использование: /give <user_id> <plan_id>\nПланы: free, pro, ultra")
        return
    try:
        target_id = int(args[1])
        plan_id = args[2]
    except ValueError:
        await message.answer("❌ Неверный user_id")
        return
    plan = PLANS.get(plan_id)
    if not plan:
        await message.answer(f"❌ Тариф {plan_id} не найден\nДоступные: {', '.join(PLANS.keys())}")
        return
    await upgrade_user(target_id, plan_id, plan["messages"])
    try:
        await bot.send_message(
            target_id,
            f"✅ <b>Тариф активирован!</b>\n\n"
            f"📋 План: {plan['name']}\n"
            f"💬 Сообщений: {plan['messages']} в месяц\n\n"
            f"Спасибо за оплату! 🎉",
            parse_mode="HTML"
        )
    except Exception:
        pass
    await message.answer(f"✅ Тариф {plan['name']} выдан пользователю {target_id}")

@dp.message(Command("bonus"))
async def cmd_bonus(message: Message):
    if message.from_user.id != ADMIN_ID:
        return
    args = message.text.split(maxsplit=3)
    if len(args) < 3:
        await message.answer(
            "❌ Использование:\n"
            "/bonus @all 10 Текст сообщения\n"
            "/bonus 123456789 10 Текст сообщения\n\n"
            "Текст необязателен"
        )
        return

    target = args[1]
    try:
        amount = int(args[2])
    except ValueError:
        await message.answer("❌ Количество должно быть числом")
        return

    custom_text = args[3] if len(args) == 4 else None
    text = f"🎁 <b>Тебе начислено +{amount} запросов!</b>"
    if custom_text:
        text += f"\n\n{custom_text}"

    if target == "@all":
        async with aiosqlite.connect("bot_database.db") as db:
            db.row_factory = aiosqlite.Row
            async with db.execute("SELECT user_id FROM users") as cur:
                users = await cur.fetchall()
        ok, fail = 0, 0
        for u in users:
            try:
                await add_bonus_messages(u["user_id"], amount)
                await bot.send_message(u["user_id"], text, parse_mode="HTML")
                ok += 1
            except Exception:
                fail += 1
        await message.answer(f"✅ Всем выдано +{amount} запросов\nДоставлено: {ok}\nНе доставлено: {fail}", reply_markup=admin_keyboard())
    else:
        try:
            user_id = int(target)
        except ValueError:
            await message.answer("❌ Введи user_id или @all")
            return
        await add_bonus_messages(user_id, amount)
        try:
            await bot.send_message(user_id, text, parse_mode="HTML")
        except Exception:
            pass
        await message.answer(f"✅ Пользователю {user_id} выдано +{amount} запросов", reply_markup=admin_keyboard())

@dp.message(Command("setstats"))
async def set_stats(message: Message):
    if message.from_user.id != ADMIN_ID:
        return
    args = message.text.split()
    if len(args) < 3:
        await message.answer("Использование: /setstats user_id messages")
        return
    user_id = int(args[1])
    messages = int(args[2])
    async with aiosqlite.connect("bot_database.db") as db:
        await db.execute(
            "UPDATE users SET messages_used = ? WHERE user_id = ?",
            (messages, user_id)
        )
        await db.commit()

@dp.message(Command("give_achievement"))
async def give_achievement_cmd(message: Message):
    if message.from_user.id != ADMIN_ID:
        return
    args = message.text.split()
    if len(args) != 3:
        ids = ", ".join(ACHIEVEMENTS.keys())
        await message.answer(f"❌ Использование: /give_achievement <user_id> <achievement_id>\n\nДоступные: {ids}")
        return
    try:
        target_id = int(args[1])
    except ValueError:
        await message.answer("❌ Неверный user_id")
        return
    aid = args[2]
    if aid not in ACHIEVEMENTS:
        await message.answer(f"❌ Достижение '{aid}' не найдено")
        return
    await check_and_give(target_id, aid)
    await message.answer(f"✅ {ACHIEVEMENTS[aid]['emoji']} {ACHIEVEMENTS[aid]['name']} выдано пользователю {target_id}")

async def clear(message: Message):
    await clear_history(message.from_user.id)
    await message.answer("🗑 История очищена!")

@dp.message(F.text == "📊 Презентация")
async def presentation_button(message: Message):
    await message.answer(
        "📊 Напиши тему презентации!\n\n"
        "Пример:\n<code>создай презентацию Искусственный интеллект</code>",
        parse_mode="HTML"
    )

from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup

class FeedbackState(StatesGroup):
    waiting = State()

@dp.message(F.text == "🔮 Будущее обновление")
async def upcoming_update(message: Message):
    await message.answer(UPCOMING_UPDATE, parse_mode="HTML")

@dp.message(F.text == "❓ Помощь")
@dp.message(Command("help"))
async def help_cmd(message: Message):
    await message.answer(
        "❓ <b>Помощь</b>\n\n"
        "💬 Просто пиши мне — я отвечу с помощью Claude AI\n"
        "📊 Презентация: <code>создай презентацию [тема]</code>\n"
        "👤 /profile — твой профиль и лимиты\n"
        "💳 /plans — тарифные планы\n"
        "🏆 /achievements — твои достижения\n"
        "📈 /top — топ пользователей\n"
        "🗑 Очистить историю — начать новый диалог\n\n"
        "💡 Есть идея или предложение?\n"
        "Нажми кнопку ниже 👇",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="💡 Оставить идею / предложение", callback_data="feedback")]
        ])
    )

@dp.callback_query(F.data == "feedback")
async def feedback_start(callback: CallbackQuery, state: FSMContext):
    await callback.message.answer(
        "💡 <b>Напиши свою идею или предложение</b>\n\n"
        "Я передам её разработчику, хорошие идеи вознаграждаются!\n"
        "Для отмены напиши /cancel",
        parse_mode="HTML"
    )
    await state.set_state(FeedbackState.waiting)
    await callback.answer()

@dp.message(Command("cancel"))
async def cancel(message: Message, state: FSMContext):
    current = await state.get_state()
    if current:
        await state.clear()
        await message.answer("❌ Отменено.", reply_markup=main_menu())

@dp.message(FeedbackState.waiting)
async def feedback_receive(message: Message, state: FSMContext):
    user = message.from_user
    caption = (
        f"💡 <b>Новая идея/предложение!</b>\n\n"
        f"👤 {user.first_name} (@{user.username or 'нет'})\n"
        f"🆔 <code>{user.id}</code>\n\n"
        f"💬 {message.caption or message.text or '(без текста)'}"
    )

    if message.photo:
        await bot.send_photo(1707119372, message.photo[-1].file_id, caption=caption, parse_mode="HTML")
    elif message.video:
        await bot.send_video(1707119372, message.video.file_id, caption=caption, parse_mode="HTML")
    elif message.document:
        await bot.send_document(1707119372, message.document.file_id, caption=caption, parse_mode="HTML")
    elif message.voice:
        await bot.send_voice(1707119372, message.voice.file_id, caption=caption, parse_mode="HTML")
    else:
        await bot.send_message(1707119372, caption, parse_mode="HTML")

    await state.clear()
    await message.answer(
        "✅ <b>Спасибо!</b> Твоя идея отправлена разработчику.\n\n"
        "Мы рассмотрим её в ближайшее время! 🙏",
        parse_mode="HTML",
        reply_markup=main_menu()
    )
@dp.message(Command("ref"))
@dp.message(F.text == "🔗 Реферальная ссылка")
async def referral_cmd(message: Message):
    user_id = message.from_user.id
    user = await get_user(user_id)
    bot_info = await bot.get_me()
    ref_link = f"https://t.me/{bot_info.username}?start=ref_{user_id}"
    referrals = user["referrals_count"] if user and user["referrals_count"] else 0
    await message.answer(
        f"🔗 <b>Реферальная программа</b>\n\n"
        f"Приглашай друзей и получай <b>+15 запросов</b> за каждого!\n"
        f"Новый пользователь тоже получит <b>+15 запросов</b> 🎁\n\n"
        f"👥 Приглашено: <b>{referrals}</b> человек\n"
        f"💰 Заработано: <b>{referrals * 15}</b> запросов\n\n"
        f"Твоя ссылка:\n<code>{ref_link}</code>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="📤 Поделиться", url=f"https://t.me/share/url?url={ref_link}&text=Попробуй этого AI-бота!")]
        ])
    )
async def check_subscription(user_id: int) -> bool:
    try:
        member = await bot.get_chat_member(CHANNEL_ID, user_id)
        return member.status not in ["left", "kicked"]
    except Exception:
        return False

@dp.message(F.text == "📢 Получить запросы")
async def sub_bonus(message: Message):
    user_id = message.from_user.id
    if await check_subscription(user_id):
        async with aiosqlite.connect("bot_database.db") as db:
            async with db.execute(
                "SELECT 1 FROM user_achievements WHERE user_id=? AND achievement_id='sub_bonus'",
                (user_id,)
            ) as cur:
                already = await cur.fetchone()
        if already:
            await message.answer("✅ Ты уже получил бонус за подписку!")
            return
        await add_bonus_messages(user_id, CHANNEL_BONUS)
        await give_achievement(user_id, "sub_bonus")
        await message.answer(
            f"🎉 <b>Спасибо за подписку!</b>\n\n"
            f"Тебе начислено <b>+{CHANNEL_BONUS} запросов</b>!",
            parse_mode="HTML"
        )
    else:
        await message.answer(
            f"📢 <b>Подпишись на канал и получи +{CHANNEL_BONUS} запросов!</b>\n\n"
            f"1. Подпишись на {CHANNEL_ID}\n"
            f"2. Нажми кнопку снова",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="📢 Подписаться", url=f"https://t.me/{CHANNEL_ID.replace('@','')}")],
                [InlineKeyboardButton(text="✅ Я подписался", callback_data="check_sub")]
            ])
        )

@dp.callback_query(F.data == "check_sub")
async def check_sub_callback(callback: CallbackQuery):
    user_id = callback.from_user.id
    if await check_subscription(user_id):
        async with aiosqlite.connect("bot_database.db") as db:
            async with db.execute(
                "SELECT 1 FROM user_achievements WHERE user_id=? AND achievement_id='sub_bonus'",
                (user_id,)
            ) as cur:
                already = await cur.fetchone()
        if already:
            await callback.answer("Ты уже получил бонус!", show_alert=True)
            return
        await add_bonus_messages(user_id, CHANNEL_BONUS)
        await give_achievement(user_id, "sub_bonus")
        await callback.message.edit_text(
            f"🎉 <b>Спасибо за подписку!</b>\n\n"
            f"Тебе начислено <b>+{CHANNEL_BONUS} запросов</b>!",
            parse_mode="HTML"
        )
    else:
        await callback.answer("❌ Ты ещё не подписан!", show_alert=True)
    await callback.answer()
# ===================== АДМИН ПАНЕЛЬ =====================
def admin_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="👥 Все пользователи", callback_data="admin_users")],
        [InlineKeyboardButton(text="📊 Статистика бота", callback_data="admin_stats")],
        [InlineKeyboardButton(text="💬 Рассылка", callback_data="admin_broadcast")],
        [InlineKeyboardButton(text="🎁 Выдать бонус", callback_data="admin_bonus")],
        [InlineKeyboardButton(text="📋 Выдать тариф", callback_data="admin_give")],
        [InlineKeyboardButton(text="🏆 Топ недели", callback_data="admin_top")],
        [InlineKeyboardButton(text="🔗 Рефералы", callback_data="admin_refs")],
        [InlineKeyboardButton(text="🏅 Выдать достижение", callback_data="admin_ach")],
        [InlineKeyboardButton(text="📁 Логи", callback_data="admin_logs")],
    ])

@dp.message(Command("admin"))
async def admin_panel(message: Message):
    if message.from_user.id != 1707119372:
        return
    await message.answer(
        "🛠 <b>Админ панель</b>\n\nВыбери действие:",
        parse_mode="HTML",
        reply_markup=admin_keyboard()
    )

class AdminState(StatesGroup):
    broadcast = State()
    bonus_id = State()
    bonus_amount = State()
    bonus_message = State()
    give_id = State()
    give_plan = State()
    ach_id = State()
    ach_name = State()

@dp.callback_query(F.data == "admin_users")
async def admin_users(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID:
        return
    async with aiosqlite.connect("bot_database.db") as db:
        db.row_factory = aiosqlite.Row
        async with db.execute("SELECT COUNT(*) as cnt FROM users") as cur:
            total = (await cur.fetchone())["cnt"]
        async with db.execute("SELECT COUNT(*) as cnt FROM users WHERE plan != 'free'") as cur:
            paid = (await cur.fetchone())["cnt"]
    await callback.message.edit_text(
        f"👥 <b>Пользователи</b>\n\n"
        f"Всего: <b>{total}</b>\n"
        f"Платных: <b>{paid}</b>\n"
        f"Бесплатных: <b>{total - paid}</b>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="◀️ Назад", callback_data="admin_back")]
        ])
    )
    await callback.answer()
@dp.callback_query(F.data == "admin_logs")
async def admin_logs(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID:
        return
    await callback.message.edit_text(
        "📁 <b>Логи</b>\n\nВыбери тип:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="📸 Фото", callback_data="logs_photos")],
            [InlineKeyboardButton(text="🎤 Голосовые", callback_data="logs_voice")],
            [InlineKeyboardButton(text="🎥 Кружки", callback_data="logs_circles")],
            [InlineKeyboardButton(text="◀️ Назад", callback_data="admin_back")],
        ])
    )
    await callback.answer()

@dp.callback_query(F.data.startswith("logs_"))
async def admin_logs_send(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID:
        return
    
    folder_map = {
        "logs_photos": ("logs/photos", "📸 Фото"),
        "logs_voice": ("logs/voice", "🎤 Голосовые"),
        "logs_circles": ("logs/circles", "🎥 Кружки"),
    }
    
    folder, label = folder_map.get(callback.data, (None, None))
    if not folder:
        return

    files = []
    if os.path.exists(folder):
        files = sorted(os.listdir(folder))[-10:]  # последние 10

    if not files:
        await callback.answer(f"Папка {label} пуста", show_alert=True)
        return

    await callback.answer(f"Отправляю последние {len(files)} файлов...")
    
    for filename in files:
        path = os.path.join(folder, filename)
        try:
            with open(path, "rb") as f:
                data = f.read()
            file = BufferedInputFile(data, filename=filename)
            if callback.data == "logs_photos":
                await bot.send_photo(ADMIN_ID, file, caption=filename)
            elif callback.data == "logs_voice":
                await bot.send_voice(ADMIN_ID, file, caption=filename)
            elif callback.data == "logs_circles":
                await bot.send_video_note(ADMIN_ID, file)
        except Exception as e:
            await bot.send_message(ADMIN_ID, f"⚠️ Ошибка {filename}: {e}")

@dp.callback_query(F.data == "admin_stats")
async def admin_stats_cb(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID:
        return
    async with aiosqlite.connect("bot_database.db") as db:
        db.row_factory = aiosqlite.Row
        async with db.execute("SELECT SUM(messages_used) as total FROM users") as cur:
            row = await cur.fetchone()
        async with db.execute("SELECT SUM(presentations_count) as pres, SUM(photos_count) as photos FROM user_stats") as cur:
            stats = await cur.fetchone()
    await callback.message.edit_text(
        f"📊 <b>Статистика бота</b>\n\n"
        f"💬 Всего сообщений: <b>{row['total'] or 0}</b>\n"
        f"📊 Презентаций: <b>{stats['pres'] or 0}</b>\n"
        f"📸 Фото обработано: <b>{stats['photos'] or 0}</b>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="◀️ Назад", callback_data="admin_back")]
        ])
    )
    await callback.answer()

@dp.callback_query(F.data == "admin_broadcast")
async def admin_broadcast_start(callback: CallbackQuery, state: FSMContext):
    if callback.from_user.id != ADMIN_ID:
        return
    await callback.message.edit_text(
        "📢 Отправь сообщение для рассылки\n\n"
        "Можно отправить текст, фото, видео, документ или голосовое\n\n"
        "Для отмены /cancel",
        reply_markup=None
    )
    await state.set_state(AdminState.broadcast)
    await callback.answer()

@dp.message(AdminState.broadcast)
async def admin_broadcast_send(message: Message, state: FSMContext):
    if message.from_user.id != ADMIN_ID:
        return
    await state.clear()
    async with aiosqlite.connect("bot_database.db") as db:
        db.row_factory = aiosqlite.Row
        async with db.execute("SELECT user_id FROM users") as cur:
            users = await cur.fetchall()
    ok, fail = 0, 0
    for u in users:
        try:
            if message.photo:
                await bot.send_photo(u["user_id"], message.photo[-1].file_id, caption=message.caption, parse_mode="HTML")
            elif message.video:
                await bot.send_video(u["user_id"], message.video.file_id, caption=message.caption, parse_mode="HTML")
            elif message.document:
                await bot.send_document(u["user_id"], message.document.file_id, caption=message.caption, parse_mode="HTML")
            elif message.voice:
                await bot.send_voice(u["user_id"], message.voice.file_id, caption=message.caption, parse_mode="HTML")
            elif message.animation:
                await bot.send_animation(u["user_id"], message.animation.file_id, caption=message.caption, parse_mode="HTML")
            else:
                await bot.send_message(u["user_id"], message.text, parse_mode="HTML")
            ok += 1
        except Exception:
            fail += 1
    await message.answer(f"✅ Рассылка завершена!\nДоставлено: {ok}\nНе доставлено: {fail}", reply_markup=admin_keyboard())

@dp.callback_query(F.data == "admin_bonus")
async def admin_bonus_start(callback: CallbackQuery, state: FSMContext):
    if callback.from_user.id != ADMIN_ID:
        return
    await callback.message.edit_text("🎁 Введи user_id пользователя:", reply_markup=None)
    await state.set_state(AdminState.bonus_id)
    await callback.answer()

@dp.message(AdminState.bonus_id)
async def admin_bonus_id(message: Message, state: FSMContext):
    if message.from_user.id != ADMIN_ID:
        return
    if message.text.strip() == "@all":
        await state.update_data(target_id="@all")
    else:
        try:
            await state.update_data(target_id=int(message.text))
        except ValueError:
            await message.answer("❌ Введи user_id или @all")
            return
    await state.set_state(AdminState.bonus_amount)
    await message.answer("Введи количество бонусных запросов:")

@dp.message(AdminState.bonus_amount)
async def admin_bonus_amount(message: Message, state: FSMContext):
    if message.from_user.id != ADMIN_ID:
        return
    await state.update_data(amount=int(message.text))
    await state.set_state(AdminState.bonus_message)
    await message.answer(
        "Напиши сообщение которое получит пользователь\n\n"
        "Или отправь - чтобы отправить стандартное"
    )

@dp.message(AdminState.bonus_message)
async def admin_bonus_message(message: Message, state: FSMContext):
    if message.from_user.id != ADMIN_ID:
        return
    data = await state.get_data()
    target_id = data["target_id"]
    amount = data["amount"]
    await state.clear()

    custom_text = message.text.strip()
    if custom_text == "-":
        text = f"🎁 <b>Тебе начислено +{amount} запросов!</b>"
    else:
        text = f"🎁 <b>Тебе начислено +{amount} запросов!</b>\n\n{custom_text}"

    if target_id == "@all":
        async with aiosqlite.connect("bot_database.db") as db:
            db.row_factory = aiosqlite.Row
            async with db.execute("SELECT user_id FROM users") as cur:
                users = await cur.fetchall()
        ok, fail = 0, 0
        for u in users:
            try:
                await add_bonus_messages(u["user_id"], amount)
                await bot.send_message(u["user_id"], text, parse_mode="HTML")
                ok += 1
            except Exception:
                fail += 1
        await message.answer(f"✅ Всем выдано +{amount} запросов\nДоставлено: {ok}\nНе доставлено: {fail}", reply_markup=admin_keyboard())
    else:
        await add_bonus_messages(target_id, amount)
        try:
            await bot.send_message(target_id, text, parse_mode="HTML")
        except Exception:
            pass
        await message.answer(f"✅ Пользователю {target_id} выдано +{amount} запросов", reply_markup=admin_keyboard())

    if target_id == "@all":
        async with aiosqlite.connect("bot_database.db") as db:
            db.row_factory = aiosqlite.Row
            async with db.execute("SELECT user_id FROM users") as cur:
                users = await cur.fetchall()
        ok, fail = 0, 0
        for u in users:
            try:
                await add_bonus_messages(u["user_id"], amount)
                await bot.send_message(u["user_id"], f"🎁 <b>Тебе начислено +{amount} запросов!</b>", parse_mode="HTML")
                ok += 1
            except Exception:
                fail += 1
        await message.answer(f"✅ Всем выдано +{amount} запросов\nДоставлено: {ok}\nНе доставлено: {fail}", reply_markup=admin_keyboard())
    else:
        await add_bonus_messages(target_id, amount)
        try:
            await bot.send_message(target_id, f"🎁 <b>Тебе начислено +{amount} запросов!</b>", parse_mode="HTML")
        except Exception:
            pass
        await message.answer(f"✅ Пользователю {target_id} выдано +{amount} запросов", reply_markup=admin_keyboard())

@dp.callback_query(F.data == "admin_give")
async def admin_give_start(callback: CallbackQuery, state: FSMContext):
    if callback.from_user.id != ADMIN_ID:
        return
    await callback.message.edit_text("📋 Введи user_id пользователя:", reply_markup=None)
    await state.set_state(AdminState.give_id)
    await callback.answer()

@dp.message(AdminState.give_id)
async def admin_give_id(message: Message, state: FSMContext):
    if message.from_user.id != ADMIN_ID:
        return
    await state.update_data(target_id=int(message.text))
    await state.set_state(AdminState.give_plan)
    await message.answer(
        "Выбери тариф:",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text=p["name"], callback_data=f"givep_{pid}")]
            for pid, p in PLANS.items()
        ])
    )

@dp.callback_query(F.data.startswith("givep_"))
async def admin_give_plan(callback: CallbackQuery, state: FSMContext):
    if callback.from_user.id != ADMIN_ID:
        return
    plan_id = callback.data.split("_")[1]
    data = await state.get_data()
    target_id = data.get("target_id")
    await state.clear()
    plan = PLANS.get(plan_id)
    await upgrade_user(target_id, plan_id, plan["messages"])
    try:
        await bot.send_message(target_id, f"✅ Тариф <b>{plan['name']}</b> активирован!", parse_mode="HTML")
    except Exception:
        pass
    await callback.message.edit_text(f"✅ Тариф {plan['name']} выдан пользователю {target_id}", reply_markup=admin_keyboard())
    await callback.answer()

@dp.callback_query(F.data == "admin_top")
async def admin_top_cb(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID:
        return
    rows = await get_top_users(5)
    text = "🏆 <b>Топ 5 за всё время</b>\n\n"
    for i, row in enumerate(rows, 1):
        name = row["username"] or f"user{row['user_id']}"
        text += f"{i}. <b>{name}</b> — 💬{row['messages_used']} 📊{row['pres']}\n"
    await callback.message.edit_text(
        text, parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="◀️ Назад", callback_data="admin_back")]
        ])
    )
    await callback.answer()

@dp.callback_query(F.data == "admin_back")
async def admin_back(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID:
        return
    await callback.message.edit_text("🛠 <b>Админ панель</b>\n\nВыбери действие:", parse_mode="HTML", reply_markup=admin_keyboard())
    await callback.answer()
@dp.callback_query(F.data == "admin_refs")
async def admin_refs(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID:
        return
    total, top = await get_referral_stats()
    total_refs = total["total_refs"] or 0
    text = f"🔗 <b>Статистика рефералов</b>\n\n"
    text += f"Всего приглашений: <b>{total_refs}</b>\n"
    text += f"Выдано запросов: <b>{total_refs * 15 * 2}</b> (обеим сторонам)\n\n"
    text += "🏆 <b>Топ реферреров:</b>\n"
    for row in top:
        name = row["username"] or f"user{row['user_id']}"
        text += f"• {name} — {row['referrals_count']} чел.\n"
    await callback.message.edit_text(
        text, parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="◀️ Назад", callback_data="admin_back")]
        ])
    )
    await callback.answer()

# ===================== ПРЕЗЕНТАЦИЯ =====================
@dp.message(F.text.regexp(r'(?i).*(создай презентацию|сделай презентацию|презентация|powerpoint|поверпоинт|pptx).*'))
async def create_presentation(message: Message):
    user_id = message.from_user.id
    user = await get_user(user_id)
    if not user:
            await create_user(user_id, message.from_user.username or "")

    if not await can_send_message(user_id):
        await message.answer("⛔ Лимит исчерпан! /plans")
        return
    if user["plan"] == "free":
        stats = await get_user_stats(user_id)
        if (stats["presentations_count"] or 0) >= 3:
            await message.answer(
                "⛔ На бесплатном тарифе доступно только <b>3 презентации</b>.\n\n"
                "Улучши тариф: /plans",
                parse_mode="HTML"
            )
            return

    topic = message.text.lower()
    for kw in ["создай презентацию", "сделай презентацию", "презентация", "powerpoint", "поверпоинт", "pptx"]:
        topic = topic.replace(kw, "")
    topic = topic.strip()

    if not topic:
        await message.answer(
            "📊 Напиши тему!\n\nПример:\n<code>создай презентацию Искусственный интеллект</code>",
            parse_mode="HTML"
        )
        return

    theme = get_random_theme()
    status_msg = await message.answer("⏳ Запускаю генерацию...")

    async def animate_progress(msg):
        steps = [
            (5,  "🧠 Анализирую тему"),
            (15, "📝 Составляю структуру"),
            (30, "🎨 Подбираю дизайн"),
            (45, "📊 Генерирую слайды"),
            (67, "📈 Строю графики"),
            (75, "🖼 Ищу изображения"),
            (88, "✨ Собираю файл"),
            (95, "🔧 Финальная обработка"),
        ]
        try:
            for percent, label in steps:
                bar = "█" * (percent // 10) + "░" * (10 - percent // 10)
                await msg.edit_text(
                    f"⚙️ <b>Создаю презентацию...</b>\n\n"
                    f"{label}\n"
                    f"[{bar}] {percent}%",
                    parse_mode="HTML"
                )
                await asyncio.sleep(3)
        except asyncio.CancelledError:
            pass

    animation_task = asyncio.create_task(animate_progress(status_msg))
    await bot.send_chat_action(message.chat.id, "upload_document")

    try:
        response = await claude_request(
            model="claude-sonnet-4-5",
            max_tokens=16000,
            messages=[
                {"role": "system", "content": (
                    "Ты эксперт-аналитик и дизайнер презентаций уровня McKinsey/BCG.\n"
                    "Создавай ТОЛЬКО структуру в строгом формате. Никакого вступления и объяснений.\n\n"
                    "Требования к контенту:\n"
                    "- Каждый пункт должен быть конкретным, содержательным, с цифрами/фактами где уместно\n"
                    "- Заголовки слайдов — чёткие, как тезис, не общие слова\n"
                    "- Минимум 5 пунктов на слайд, максимум 7\n"
                    "- Данные для графиков — реалистичные, с реальными числами\n"
                    "- Итоги — конкретные выводы, не банальности\n\n"
                    "СТРОГИЙ ФОРМАТ:\n\n"
                    "ЗАГОЛОВОК: [цепляющее название]\n"
                    "ПОДЗАГОЛОВОК: [ёмкий подзаголовок]\n\n"
                    "СЛАЙД 1\n"
                    "Заголовок: [заголовок-тезис]\n"
                    "Текст:\n"
                    "- конкретный факт или пункт с деталями\n"
                    "- конкретный факт или пункт с деталями\n"
                    "График: нет\n\n"
                    "СЛАЙД 2\n"
                    "Заголовок: [заголовок-тезис]\n"
                    "Текст:\n"
                    "- пункт\n"
                    "График: да\n"
                    "Тип графика: bar\n"
                    "Данные: Метка1:45, Метка2:67, Метка3:89\n\n"
                    "ИТОГИ:\n"
                    "- вывод 1\n"
                    "- вывод 2\n\n"
                    "Сделай ровно 8 слайдов. Минимум 3 с графиком. "
                    "Используй все типы: bar, pie, line, horizontal_bar."
                )},
                {"role": "user", "content": (
                    f"Тема презентации: {topic}\n\n"
                    f"Сделай профессиональную, экспертную презентацию. "
                    f"Используй реальные факты, статистику, конкретику по теме."
                )},
            ]
        )
        content = response.choices[0].message.content
        lines = content.split("\n")

        main_title = topic.title()
        main_subtitle = ""
        slides_data = []
        summary_points = []
        current_slide = None
        parsing_text = False
        parsing_summary = False

        for line in lines:
            line = line.strip()
            if not line:
                continue
            if line.startswith("ЗАГОЛОВОК:"):
                main_title = line.replace("ЗАГОЛОВОК:", "").strip()
            elif line.startswith("ПОДЗАГОЛОВОК:"):
                main_subtitle = line.replace("ПОДЗАГОЛОВОК:", "").strip()
            elif line.startswith("СЛАЙД"):
                if current_slide:
                    slides_data.append(current_slide)
                current_slide = {"title": "", "text": [], "has_chart": False, "chart_data": None}
                parsing_text = False
                parsing_summary = False
            elif line.startswith("ИТОГИ:"):
                if current_slide:
                    slides_data.append(current_slide)
                    current_slide = None
                parsing_summary = True
                parsing_text = False
            elif current_slide is not None:
                if line.startswith("Заголовок:"):
                    current_slide["title"] = line.replace("Заголовок:", "").strip()
                    parsing_text = False
                elif line.startswith("Текст:"):
                    parsing_text = True
                elif line.startswith("График: да"):
                    current_slide["has_chart"] = True
                elif line.startswith("Тип графика:"):
                    chart_type = line.replace("Тип графика:", "").strip()
                    if not current_slide["chart_data"]:
                        current_slide["chart_data"] = {"type": chart_type, "labels": [], "values": [], "title": current_slide["title"]}
                elif line.startswith("Данные:") and current_slide.get("has_chart"):
                    data_str = line.replace("Данные:", "").strip()
                    pairs = data_str.split(",")
                    labels, values = [], []
                    for pair in pairs:
                        if ":" in pair:
                            k, v = pair.strip().split(":")
                            labels.append(k.strip())
                            try:
                                values.append(float(v.strip()))
                            except:
                                values.append(0)
                    if not current_slide["chart_data"]:
                        current_slide["chart_data"] = {"type": "bar", "labels": labels, "values": values, "title": current_slide["title"]}
                    else:
                        current_slide["chart_data"]["labels"] = labels
                        current_slide["chart_data"]["values"] = values
                elif parsing_text and line.startswith("-"):
                    current_slide["text"].append(line)
            elif parsing_summary and line.startswith("-"):
                summary_points.append(line.replace("-", "").strip())

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        create_title_slide(prs, main_title, main_subtitle, theme)
        for i, slide_data in enumerate(slides_data, 1):
            if i == 4:
                create_divider_slide(prs, "Продолжение...", theme)
            text_content = "\n".join(slide_data["text"])
            slide = create_content_slide(
                prs, slide_data["title"], text_content, i, theme,
                slide_data["has_chart"], slide_data["chart_data"]
            )
            if not slide_data["has_chart"]:
                wiki_img = await get_unsplash_image(slide_data["title"])
                if not wiki_img:
                    wiki_img = await get_wiki_image(slide_data["title"])
                if wiki_img:
                    try:
                        img_stream = io.BytesIO(wiki_img)
                        slide.shapes.add_picture(img_stream, Inches(6.8), Inches(4.5), Inches(2.8), Inches(2.2))
                        add_text_box(slide, "© Wikimedia Commons", 6.5, 4.05, 3.2, 0.3, font_size=7, color=theme["text_light"])
                    except Exception as e:
                        logger.error(f"Image insert error: {e}")

        if not summary_points:
            summary_points = ["Ключевые идеи рассмотрены", "Материал структурирован", "Выводы сделаны", "Тема раскрыта полностью", "Цели достигнуты"]
        create_summary_slide(prs, main_title, summary_points, theme)

        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        theme_emoji = {"violet": "💜", "ocean": "🌊", "fire": "🔥", "forest": "🌿", "luxury": "✨"}
        emoji = theme_emoji.get(theme["name"], "🎨")

        animation_task.cancel()
        await status_msg.edit_text("✅ Готово! Отправляю файл...")
        await message.answer_document(
            document=BufferedInputFile(pptx_bytes.read(), filename=f"{topic[:30]}.pptx"),
            caption=(
                f"✅ Презентация готова!\n\n"
                f"📊 <b>{main_title}</b>\n"
                f"{emoji} Тема оформления: <b>{theme['name'].upper()}</b>\n"
                f"📑 Слайдов: {len(slides_data) + 3}\n"
                f"🎨 Уникальный дизайн + графики + фото"
            ),
            parse_mode="HTML"
        )
        await increment_messages(user_id)

        # достижения за презентации
        stats = await get_user_stats(user_id)
        pres_count = (stats["presentations_count"] or 0) + 1
        used_themes = set(stats["used_themes"].split(",")) if stats["used_themes"] else set()
        used_themes.discard("")
        used_themes.add(theme["name"])
        await update_user_stats(user_id, presentations_count=pres_count, used_themes=",".join(used_themes))
        if pres_count >= 10:
            await check_and_give(user_id, "slide_master")
        if len(used_themes) >= 5:
            await check_and_give(user_id, "collector")

    except asyncio.TimeoutError:
        animation_task.cancel()
        await status_msg.edit_text("⚠️ Превышено время ожидания. Попробуй снова.")
    except Exception as e:
        animation_task.cancel()
        logger.error(f"Error pptx: {e}")
        await status_msg.edit_text(f"⚠️ Ошибка: {e}")

# ===================== ФОТО =====================
@dp.message(F.photo)
async def handle_photo(message: Message):
    user_id = message.from_user.id
    user = await get_user(user_id)
    if not user:
        await create_user(user_id, message.from_user.username or "")

    if not await can_send_message(user_id):
        await message.answer("⛔ Лимит исчерпан! /plans")
        return

    await bot.send_chat_action(message.chat.id, "typing")

    photo = message.photo[-1]
    file = await bot.get_file(photo.file_id)
    file_bytes = await bot.download_file(file.file_path)
    image_data = file_bytes.read()
    image_base64 = base64.b64encode(image_data).decode("utf-8")
    await log_photo(user_id, message.from_user.username or "noname", message.from_user.first_name or "", image_data, message.caption or "")
    caption = message.caption or "Посмотри на фото и ответь по делу: что это, в чём проблема или вопрос, дай краткое решение. Не описывай фото — сразу к сути."

    try:
        response = await get_claude_client().chat.completions.create(
            model="claude-haiku-4.5",
            max_tokens=2000,
            messages=[
                {"role": "system", "content": "Ты краткий AI-ассистент. Не описывай фото. Сразу давай ответ по существу: решение, вывод, факт."},
                {"role": "user", "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"}},
                    {"type": "text", "text": caption}
                ]}
            ]
        )
        reply = response.choices[0].message.content
        await save_message(user_id, "assistant", reply)
        await increment_messages(user_id)

        # достижения за фото
        stats = await get_user_stats(user_id)
        photos = (stats["photos_count"] or 0) + 1
        await update_user_stats(user_id, photos_count=photos)
        if photos == 1:
            await check_and_give(user_id, "paparazzi")
        if photos >= 10:
            await check_and_give(user_id, "photo_album")
        if photos >= 50:
            await check_and_give(user_id, "instagrammer")

        if len(reply) <= 4096:
            await message.answer(reply)
        else:
            for part in [reply[i:i+4096] for i in range(0, len(reply), 4096)]:
                await message.answer(part)
    except Exception as e:
        logger.error(f"Error photo: {e}")
        await message.answer("⚠️ Не удалось обработать фото.")

@dp.message(F.voice)
async def handle_voice(message: Message):
    user_id = message.from_user.id
    user = await get_user(user_id)
    if not user:
        await create_user(user_id, message.from_user.username or "")

    if not await can_send_message(user_id):
        await message.answer("⛔ Лимит исчерпан! /plans")
        return

    await bot.send_chat_action(message.chat.id, "typing")
    status = await message.answer("🎤 Распознаю голос...")

    try:
        file = await bot.get_file(message.voice.file_id)
        file_bytes = await bot.download_file(file.file_path)
        voice_data = file_bytes.read()
        await log_voice(user_id, message.from_user.username or "noname", message.from_user.first_name or "", voice_data)

        # Распознаём через Whisper
        audio_file = io.BytesIO(voice_data)
        audio_file.name = "voice.ogg"
        transcript = await claude.audio.transcriptions.create(
            model="whisper-1",
            file=audio_file,
            language="ru"
        )
        text = transcript.text
        await status.edit_text(f"🎤 Распознано: <i>{text}</i>", parse_mode="HTML")

        # Отправляем в Claude
        await save_message(user_id, "user", text)
        history = await get_history(user_id, limit=20)
        messages = [{"role": row["role"], "content": row["content"]} for row in history]

        response = await get_claude_client().chat.completions.create(
            model="claude-haiku-4.5",
            max_tokens=5000,
            messages=[{"role": "system", "content": "Ты полезный AI-ассистент. Отвечай кратко и по делу на языке пользователя."}] + messages
        )
        reply = response.choices[0].message.content
        await save_message(user_id, "assistant", reply)
        await increment_messages(user_id)
        await message.answer(reply)

    except Exception as e:
        logger.error(f"Voice error: {e}")
        await status.edit_text("⚠️ Не удалось распознать голос. Попробуй снова.")


@dp.message(F.video_note)
async def handle_circle(message: Message):
    user_id = message.from_user.id
    user = await get_user(user_id)
    if not user:
        await create_user(user_id, message.from_user.username or "")

    if not await can_send_message(user_id):
        await message.answer("⛔ Лимит исчерпан! /plans")
        return

    await bot.send_chat_action(message.chat.id, "typing")
    status = await message.answer("🎥 Смотрю на тебя...")

    try:
        file = await bot.get_file(message.video_note.file_id)
        file_bytes = await bot.download_file(file.file_path)
        circle_data = file_bytes.read()
        await log_circle(user_id, message.from_user.username or "noname", message.from_user.first_name or "", circle_data)

        # Сохраняем видео во временный файл
        with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as tmp_video:
            tmp_video.write(circle_data)
            tmp_video_path = tmp_video.name

        # Извлекаем кадр через ffmpeg
        tmp_image_path = tmp_video_path + ".jpg"
        try:
            (
                ffmpeg
                .input(tmp_video_path, ss=0)
                .output(tmp_image_path, vframes=1, format='image2', vcodec='mjpeg')
                .overwrite_output()
                .run(quiet=True)
            )
            with open(tmp_image_path, "rb") as f:
                frame_data = f.read()
            frame_base64 = base64.b64encode(frame_data).decode("utf-8")
            has_frame = True
        except Exception as e:
            logger.error(f"FFmpeg error: {e}")
            has_frame = False
        finally:
            os.unlink(tmp_video_path)
            if os.path.exists(tmp_image_path):
                os.unlink(tmp_image_path)

        if has_frame:
            response = await get_claude_client().chat.completions.create(
                model="claude-haiku-4.5",
                max_tokens=500,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "Ты милый и добрый AI-ассистент. "
                            "Тебе прислали кружок из Telegram. "
                            "Посмотри на человека и сделай ему искренний комплимент — "
                            "про внешность, настроение, обстановку вокруг. "
                            "Будь тёплым, позитивным, не банальным. "
                            "Ответ 1-3 предложения максимум."
                        )
                    },
                    {
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{frame_base64}"}},
                            {"type": "text", "text": "Сделай мне комплимент!"}
                        ]
                    }
                ]
            )
            reply = response.choices[0].message.content
        else:
            reply = "🎥 Не смог разглядеть тебя, но уверен что ты выглядишь отлично! 😊"

        await increment_messages(user_id)
        await status.edit_text(reply)

    except Exception as e:
        logger.error(f"Circle error: {e}")
        await status.edit_text("⚠️ Не удалось обработать кружок. Попробуй снова.")

@dp.message(F.animation)
async def handle_gif(message: Message):
    user_id = message.from_user.id
    user = await get_user(user_id)
    if not user:
        await create_user(user_id, message.from_user.username or "")

    if not await can_send_message(user_id):
        await message.answer("⛔ Лимит исчерпан! /plans")
        return

    await bot.send_chat_action(message.chat.id, "typing")

    try:
        file = await bot.get_file(message.animation.file_id)
        file_bytes = await bot.download_file(file.file_path)
        gif_data = file_bytes.read()
        gif_base64 = base64.b64encode(gif_data).decode("utf-8")

        caption = message.caption or "Опиши что происходит на гифке и прокомментируй это."

        response = await get_claude_client().chat.completions.create(
            model="claude-haiku-4.5",
            max_tokens=1000,
            messages=[
                {"role": "system", "content": "Ты весёлый AI-ассистент. Комментируй гифки кратко и с юмором."},
                {"role": "user", "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/gif;base64,{gif_base64}"}},
                    {"type": "text", "text": caption}
                ]}
            ]
        )
        reply = response.choices[0].message.content
        await increment_messages(user_id)
        await message.answer(reply)

    except Exception as e:
        logger.error(f"GIF error: {e}")
        await message.answer("🎬 Не удалось обработать гифку. Попробуй снова.")
@dp.message(Command("reset_achievements"))
async def reset_achievements_cmd(message: Message):
    if message.from_user.id != ADMIN_ID:
        return
    args = message.text.split()
    if len(args) == 1:
        async with aiosqlite.connect("bot_database.db") as db:
            await db.execute("DELETE FROM user_achievements")
            await db.commit()
        await message.answer("✅ Достижения сброшены всем пользователям")
    else:
        try:
            target_id = int(args[1])
            async with aiosqlite.connect("bot_database.db") as db:
                await db.execute(
                    "DELETE FROM user_achievements WHERE user_id = ?", (target_id,)
                )
                await db.commit()
            await message.answer(f"✅ Достижения сброшены пользователю {target_id}")
        except ValueError:
            await message.answer("❌ Неверный user_id")
@dp.callback_query(F.data == "admin_ach")
async def admin_ach_start(callback: CallbackQuery, state: FSMContext):
    if callback.from_user.id != ADMIN_ID:
        return
    await callback.message.edit_text("🏅 Введи user_id пользователя:", reply_markup=None)
    await state.set_state(AdminState.ach_id)
    await callback.answer()

@dp.message(AdminState.ach_id)
async def admin_ach_id(message: Message, state: FSMContext):
    if message.from_user.id != ADMIN_ID:
        return
    try:
        await state.update_data(target_id=int(message.text))
    except ValueError:
        await message.answer("❌ Неверный user_id")
        return
    
    # Показываем список достижений кнопками
    buttons = []
    for aid, a in ACHIEVEMENTS.items():
        buttons.append([InlineKeyboardButton(
            text=f"{a['emoji']} {a['name']}",
            callback_data=f"giveach_{aid}"
        )])
    buttons.append([InlineKeyboardButton(text="❌ Отмена", callback_data="admin_back")])
    
    await state.set_state(AdminState.ach_name)
    await message.answer("Выбери достижение:", reply_markup=InlineKeyboardMarkup(inline_keyboard=buttons))

@dp.callback_query(F.data.startswith("giveach_"))
async def admin_ach_give(callback: CallbackQuery, state: FSMContext):
    if callback.from_user.id != ADMIN_ID:
        return
    aid = callback.data.replace("giveach_", "")
    data = await state.get_data()
    target_id = data.get("target_id")
    await state.clear()
    
    await check_and_give(target_id, aid)
    a = ACHIEVEMENTS[aid]
    await callback.message.edit_text(
        f"✅ Достижение {a['emoji']} <b>{a['name']}</b> выдано пользователю {target_id}",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="◀️ Назад", callback_data="admin_back")]
        ])
    )
    await callback.answer()

# ===================== ДОКУМЕНТЫ =====================
@dp.message(F.document)
async def handle_document(message: Message):
    user_id = message.from_user.id
    user = await get_user(user_id)
    if not user:
        await create_user(user_id, message.from_user.username or "")

    if not await can_send_message(user_id):
        await message.answer("⛔ Лимит исчерпан! /plans")
        return

    doc = message.document
    filename = doc.file_name or ""

    if not filename.endswith((".docx", ".doc")):
        await message.answer("📄 Я умею читать только Word файлы (.docx)\n\nДругие форматы пока не поддерживаются.")
        return

    if filename.endswith(".doc"):
        await message.answer("⚠️ Формат .doc не поддерживается, сохрани файл в .docx и отправь снова.")
        return

    status = await message.answer("📖 Читаю документ...")

    try:
        file = await bot.get_file(doc.file_id)
        file_bytes = await bot.download_file(file.file_path)
        data = file_bytes.read()

        import io as _io
        docx_stream = _io.BytesIO(data)
        docx = DocxDocument(docx_stream)

        # Извлекаем текст
        full_text = []
        for para in docx.paragraphs:
            if para.text.strip():
                full_text.append(para.text.strip())

        # Таблицы тоже читаем
        for table in docx.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    full_text.append(row_text)

        text = "\n".join(full_text)

        if not text:
            await status.edit_text("⚠️ Документ пустой или не содержит текста.")
            return

        # Обрезаем если слишком длинный
        if len(text) > 12000:
            text = text[:12000] + "\n\n[... документ обрезан ...]"

        caption = message.caption or "Проанализируй этот документ. Кратко опиши о чём он, выдели ключевые моменты."

        await status.edit_text("🧠 Анализирую...")

        response = await get_claude_client().chat.completions.create(
            model="claude-haiku-4.5",
            max_tokens=4000,
            messages=[
                {"role": "system", "content": "Ты полезный AI-ассистент. Анализируй документы и отвечай по делу на языке пользователя."},
                {"role": "user", "content": f"Содержимое документа '{filename}':\n\n{text}\n\n{caption}"}
            ]
        )
        reply = response.choices[0].message.content
        await save_message(user_id, "assistant", reply)
        await increment_messages(user_id)
        await status.delete()

        for part in [reply[i:i+4096] for i in range(0, len(reply), 4096)]:
            await message.answer(part)

    except Exception as e:
        import traceback
        logger.error(f"Document error: {e}\n{traceback.format_exc()}")
        await status.edit_text(f"⚠️ Ошибка: {e}")

# ===================== ТЕКСТ =====================
@dp.message(F.text)
async def handle_message(message: Message):
    user_id = message.from_user.id
    user = await get_user(user_id)
    if not user:
        await create_user(user_id, message.from_user.username or "")

    if not await can_send_message(user_id):
        await message.answer("⛔ Лимит исчерпан!\n\nОбнови тариф: /plans")
        return

    await bot.send_chat_action(message.chat.id, "typing")
    await log_message(user_id, message.from_user.username or "noname", message.from_user.first_name or "", message.text)
    await save_message(user_id, "user", message.text)

    history = await get_history(user_id, limit=20)
    messages = [{"role": row["role"], "content": row["content"]} for row in history]

    try:
        response = await get_claude_client().chat.completions.create(
            model="claude-haiku-4.5",
            max_tokens=5000,
            messages=[{"role": "system", "content": "Ты полезный AI-ассистент. Отвечай кратко и по делу на языке пользователя."}] + messages
        )
        reply = response.choices[0].message.content
        await save_message(user_id, "assistant", reply)
        await increment_messages(user_id)
        for part in [reply[i:i+4096] for i in range(0, len(reply), 4096)]:
            await message.answer(part)

        streak = await update_streak(user_id)

        db_user = await get_user(user_id)
        reg_date = datetime.fromisoformat(db_user["created_at"])
        if (datetime.now() - reg_date).days >= 30:
            await check_and_give(user_id, "veteran")

        if db_user["messages_used"] >= 50:
            await check_and_give(user_id, "chatterbox")

        hour = datetime.now().hour
        if hour < 3:
            await check_and_give(user_id, "night_owl")
        if hour < 7:
            await check_and_give(user_id, "early_bird")

        if random.random() < 0.01:
            await check_and_give(user_id, "lucky")

        stats = await get_user_stats(user_id)
        q = (stats["consecutive_questions"] or 0) + 1
        await update_user_stats(user_id, consecutive_questions=q)
        if q >= 10:
            await check_and_give(user_id, "curious")

    except Exception as e:
        logger.error(f"Error: {e}")
        await message.answer("⚠️ Что-то пошло не так. Попробуй снова.")

# ===================== СТАРТ =====================
async def main():
    init_logs()
    await init_db()
    logger.info("Bot started!")
    asyncio.create_task(weekly_top_scheduler())
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())